"""
sopno/tools/builtins/readers.py
━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
Binary document readers for Sopno — PDF, images, and Office files.

Every extractor is optional and imported lazily, so Sopno works fine without
PyMuPDF, Tesseract, Pillow, or the Office parsers. ``extract_text`` routes by
file suffix through a layered pipeline (fast native text extraction first, OCR
fallback for scanned pages/images when ``file_ocr_enabled``), caps pages and
output, and reports which path it took:

    (text, method)   e.g. ("Invoice No. 42", "pdf-native")
                     or  ("...", "pdf-ocr") / ("...", "image-ocr")
"""

from __future__ import annotations

import io
import struct
from pathlib import Path
from typing import Optional

from sopno.config.settings import settings

# Suffixes routed through the readers (everything else is plain text).
_PDF_SUFFIXES = {".pdf"}
_IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tif", ".tiff"}
_OFFICE_SUFFIXES = {".docx", ".xlsx", ".pptx"}
# Legacy Office formats — extracted via LibreOffice headless conversion.
_LEGACY_SUFFIXES = {".doc", ".xls", ".ppt", ".odt", ".ods", ".odp"}


def _max_pages() -> int:
    return max(1, int(getattr(settings, "readers_max_pages", 20)))


def _max_chars() -> int:
    return max(1, int(getattr(settings, "readers_max_chars", 20000)))


def _ocr_enabled() -> bool:
    return bool(getattr(settings, "file_ocr_enabled", True))


def _cap(text: str, method: str) -> tuple[str, str]:
    limit = _max_chars()
    if len(text) > limit:
        text = text[:limit] + "\n…(content truncated)…"
    return text or "(no readable text)", method


def _try_tesseract(path: Path) -> Optional[str]:
    """OCR a page/image with Tesseract; None if unavailable or it fails."""
    if not _ocr_enabled():
        return None
    try:
        import pytesseract
        from PIL import Image
    except Exception:  # noqa: BLE001
        return None
    try:
        return pytesseract.image_to_string(Image.open(path))
    except Exception:  # noqa: BLE001
        return None


def _ocr_image_bytes(data: bytes) -> Optional[str]:
    """OCR in-memory image bytes; None if Tesseract is unavailable/fails."""
    if not _ocr_enabled():
        return None
    try:
        import pytesseract
        from PIL import Image
    except Exception:  # noqa: BLE001
        return None
    try:
        return pytesseract.image_to_string(Image.open(io.BytesIO(data)))
    except Exception:  # noqa: BLE001
        return None


def _read_pdf(path: Path, max_pages: int) -> tuple[str, str]:
    """PDF text via PyMuPDF; OCR fallback for scanned pages."""
    try:
        import pymupdf  # type: ignore[import-not-found]
    except Exception:  # noqa: BLE001 — PyMuPDF not installed
        return "Could not read PDF — pymupdf is not installed (pip install pymupdf).", "pdf-unavailable"
    try:
        doc = pymupdf.open(path)
    except Exception as e:  # noqa: BLE001
        return f"Could not open PDF {path}: {e}", "pdf-unavailable"

    pages: list[str] = []
    ocr_pages: list[str] = []
    total = len(doc)
    for i, page in enumerate(doc):
        if i >= max_pages:
            pages.append(f"…({max(0, total - max_pages)} more pages skipped)…")
            break
        text = page.get_text().strip()
        if text:
            pages.append(text)
        elif _ocr_enabled():
            try:
                pix = page.get_pixmap(dpi=200)
                ocr = _ocr_image_bytes(pix.tobytes("png"))
                if ocr and ocr.strip():
                    ocr_pages.append(ocr.strip())
            except Exception:  # noqa: BLE001
                pass
    doc.close()

    if pages:
        return _cap("\n\n".join(pages), "pdf-native")
    if ocr_pages:
        return _cap("\n\n".join(ocr_pages), "pdf-ocr")
    if not _ocr_enabled():
        return "(no readable text — scanned PDF, and OCR is disabled)", "pdf-ocr-unavailable"
    return "(no readable text — scanned PDF without OCR support)", "pdf-ocr-unavailable"


def _read_image(path: Path) -> tuple[str, str]:
    """Image OCR via Tesseract (requires tesseract-ocr on the system)."""
    text = _ocr_image_bytes(path.read_bytes()) if _ocr_enabled() else None
    if text is None:
        return ("Image OCR unavailable — install tesseract-ocr and Pillow "
                "(pip install pillow pytesseract)."), "image-unavailable"
    return _cap(text, "image-ocr")


def _read_docx(path: Path) -> tuple[str, str]:
    try:
        import docx  # type: ignore[import-not-found]
    except Exception:  # noqa: BLE001
        return "Could not read .docx — python-docx is not installed.", "docx-unavailable"
    try:
        d = docx.Document(path)
        paras = [p.text for p in d.paragraphs if p.text.strip()]
        for table in d.tables:
            for row in table.rows:
                paras.append(" | ".join(c.text.strip() for c in row.cells))
        return _cap("\n".join(paras), "docx-native")
    except Exception as e:  # noqa: BLE001
        return f"Could not read .docx {path}: {e}", "docx-unavailable"


def _read_pptx(path: Path) -> tuple[str, str]:
    try:
        from pptx import Presentation  # type: ignore[import-not-found]
    except Exception:  # noqa: BLE001
        return "Could not read .pptx — python-pptx is not installed.", "pptx-unavailable"
    try:
        prs = Presentation(path)
        lines: list[str] = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    for para in shape.text_frame.paragraphs:
                        t = "".join(run.text for run in para.runs).strip()
                        if t:
                            lines.append(t)
        return _cap("\n".join(lines), "pptx-native")
    except Exception as e:  # noqa: BLE001
        return f"Could not read .pptx {path}: {e}", "pptx-unavailable"


def _read_xlsx(path: Path) -> tuple[str, str]:
    try:
        import openpyxl  # type: ignore[import-not-found]
    except Exception:  # noqa: BLE001
        return "Could not read .xlsx — openpyxl is not installed.", "xlsx-unavailable"
    try:
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        lines: list[str] = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                values = [str(c).strip() for c in row if c is not None and str(c).strip()]
                if values:
                    lines.append(" | ".join(values))
        wb.close()
        return _cap("\n".join(lines), "xlsx-native")
    except Exception as e:  # noqa: BLE001
        return f"Could not read .xlsx {path}: {e}", "xlsx-unavailable"


def _read_legacy(path: Path) -> tuple[str, str]:
    """Legacy Office docs (.doc/.xls/.ppt/…) via LibreOffice headless."""
    import subprocess
    import tempfile
    try:
        with tempfile.TemporaryDirectory(prefix="sopno-lo-") as td:
            outdir = Path(td)
            subprocess.run(
                [
                    "libreoffice", "--headless", "--convert-to", "txt:Text",
                    "--outdir", str(outdir), str(path),
                ],
                capture_output=True,
                timeout=120,
                check=False,
            )
            converted = outdir / (path.stem + ".txt")
            if converted.is_file():
                text = converted.read_text(encoding="utf-8", errors="replace")
                return _cap(text, "legacy-libreoffice")
    except FileNotFoundError:
        return "Could not read legacy format — install libreoffice.", "legacy-unavailable"
    except Exception:  # noqa: BLE001
        pass
    return "Could not read legacy document.", "legacy-unavailable"


def extract_text(path: Path) -> tuple[str, str]:
    """
    Extract readable text from a binary document.

    Returns ``(text, method)``. For plain text files callers should read the
    file directly; this returns ``("", "")`` when the suffix is not ours.
    """
    suffix = path.suffix.lower()
    if suffix in _PDF_SUFFIXES:
        return _read_pdf(path, _max_pages())
    if suffix in _IMAGE_SUFFIXES:
        return _read_image(path)
    if suffix in _OFFICE_SUFFIXES:
        if suffix == ".docx":
            return _read_docx(path)
        if suffix == ".pptx":
            return _read_pptx(path)
        return _read_xlsx(path)
    if suffix in _LEGACY_SUFFIXES:
        return _read_legacy(path)
    return "", ""


def is_binary_like(path: Path) -> bool:
    """True for document suffixes the readers handle (or a null-byte file)."""
    if path.suffix.lower() in (_PDF_SUFFIXES | _IMAGE_SUFFIXES | _OFFICE_SUFFIXES | _LEGACY_SUFFIXES):
        return True
    try:
        with open(path, "rb") as fh:
            head = fh.read(1024)
        return b"\x00" in head
    except OSError:
        return False
